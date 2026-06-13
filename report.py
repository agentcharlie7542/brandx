"""리포트 생성기 — 포지셔닝 맵 · 디자인 갤러리 · 갭 분석 · PPTX.

사용:
  python report.py --group skincare          # 맵(png) + 갭표(csv) 생성
  python report.py --group skincare --pptx    # + PPTX 보고서
  python report.py --pptx --merge             # 전 그룹 통합 1부

산출물(out/):
  positioning_{group}_price_review.png  — 가격(중앙값) × 리뷰볼륨
  positioning_{group}_follower_sat.png  — 팔로워 × 만족도
  gap_{group}.csv                       — 지표별 (자사 − 그룹 1위)
  report_{group}.pptx | report_merged.pptx

PPTX 구성: ①표지 ②요약 ③스코어카드 ④포지셔닝 맵 ⑤갤러리 ⑥갭→개선과제 ⑦부록.
"""
import sys
from pathlib import Path

import pandas as pd

from db import connect
from score import build_scorecard
from shops_groups import GROUPS, group_label, own_shop_ids, resolve_group

ROOT = Path(__file__).parent
OUT = ROOT / "out"

# 헤드라인 지표(higher-is-better) — 갭 분석 대상
HEADLINE = {
    "review_volume": "리뷰 볼륨",
    "est_revenue_30d": "추정 매출(30일)",
    "design_total": "디자인 총점",
    "follower_count": "팔로워",
    "satisfaction_pct": "만족도%",
    "promo_intensity": "프로모션 강도",
}


def _arg(flag: str, default=None):
    return sys.argv[sys.argv.index(flag) + 1] if flag in sys.argv else default


def _setup_matplotlib():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    # 한글 라벨 대비 CJK 폰트 시도(macOS), 실패 시 기본
    for fam in ("AppleGothic", "Apple SD Gothic Neo", "NanumGothic", "Malgun Gothic"):
        try:
            matplotlib.font_manager.findfont(fam, fallback_to_default=False)
            plt.rcParams["font.family"] = fam
            break
        except Exception:
            continue
    plt.rcParams["axes.unicode_minus"] = False
    return plt


# ── 포지셔닝 맵 ─────────────────────────────────────────────────────────────
def _scatter(df, xcol, ycol, title, xlabel, ylabel, own, path):
    plt = _setup_matplotlib()
    sub = df.dropna(subset=[xcol, ycol])
    if sub.empty:
        return None
    fig, ax = plt.subplots(figsize=(8, 6))
    for _, r in sub.iterrows():
        is_own = r["shop_id"] in own
        ax.scatter(r[xcol], r[ycol],
                   s=320 if is_own else 160,
                   c="#C00000" if is_own else "#8FAADC",
                   marker="*" if is_own else "o",
                   edgecolors="black", linewidths=0.7, zorder=3)
        ax.annotate(r["shop_id"], (r[xcol], r[ycol]),
                    fontsize=9, fontweight="bold" if is_own else "normal",
                    xytext=(6, 6), textcoords="offset points")
    # 사분면 기준선(중앙값)
    ax.axvline(sub[xcol].median(), color="gray", ls="--", lw=0.6, alpha=0.6)
    ax.axhline(sub[ycol].median(), color="gray", ls="--", lw=0.6, alpha=0.6)
    ax.set_title(title, fontsize=13, fontweight="bold")
    ax.set_xlabel(xlabel)
    ax.set_ylabel(ylabel)
    ax.grid(True, alpha=0.25)
    fig.tight_layout()
    fig.savefig(path, dpi=130)
    plt.close(fig)
    return path


def positioning_maps(df, group, own) -> list[Path]:
    OUT.mkdir(exist_ok=True)
    paths = []
    p1 = _scatter(df, "price_median", "review_volume",
                  f"포지셔닝: 가격 × 리뷰볼륨 ({group_label(group)})",
                  "가격 중앙값 (¥)", "리뷰 볼륨 (Σ)", own,
                  OUT / f"positioning_{group}_price_review.png")
    if p1:
        paths.append(p1)
    p2 = _scatter(df, "follower_count", "satisfaction_pct",
                  f"포지셔닝: 팔로워 × 만족도 ({group_label(group)})",
                  "팔로워 수", "만족도 (%)", own,
                  OUT / f"positioning_{group}_follower_sat.png")
    if p2:
        paths.append(p2)
    return paths


# ── 갭 분석 ─────────────────────────────────────────────────────────────────
def gap_table(df, own) -> pd.DataFrame:
    """지표별 (대표 자사 − 그룹 1위). 갭 오름차순 = 개선 우선순위."""
    if df.empty:
        return pd.DataFrame()
    own_df = df[df["shop_id"].isin(own)]
    if own_df.empty:
        return pd.DataFrame()
    # 대표 자사 = 리뷰볼륨 최대 자사 샵
    rep = (own_df.sort_values("review_volume", ascending=False).iloc[0]
           if "review_volume" in own_df else own_df.iloc[0])
    rows = []
    for metric, label in HEADLINE.items():
        if metric not in df or not df[metric].notna().any():
            continue
        leader_val = df[metric].max()
        leader = df.loc[df[metric].idxmax(), "shop_id"]
        own_val = rep.get(metric)
        if pd.isna(own_val):
            continue
        rows.append({
            "지표": label,
            "자사": rep["shop_id"],
            "자사값": round(float(own_val), 2),
            "그룹1위": leader,
            "1위값": round(float(leader_val), 2),
            "갭": round(float(own_val) - float(leader_val), 2),
            "달성률%": round(float(own_val) / float(leader_val) * 100, 1) if leader_val else None,
        })
    out = pd.DataFrame(rows)
    return out.sort_values("갭").reset_index(drop=True) if not out.empty else out


# ── 디자인 갤러리 (샵별 대표 썸네일) ────────────────────────────────────────
def gallery_assets(conn, group, own) -> dict:
    """그룹 Best(디자인 1위) vs 자사 대표 썸네일 경로."""
    shop_ids = [s["shop_id"] for s in resolve_group(group)]
    if not shop_ids:
        return {}
    rows = conn.execute(
        f"""SELECT shop_id, path FROM image_assets
            WHERE asset_type='thumb' AND shop_id IN ({",".join("?" * len(shop_ids))})
            ORDER BY shop_id, captured_at""", shop_ids).fetchall()
    rep = {}
    for shop_id, path in rows:
        rep.setdefault(shop_id, path)  # 샵별 첫 썸네일
    return rep


# ── PPTX ───────────────────────────────────────────────────────────────────
def _new_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)
    return prs


def _title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def _section_title(prs, text):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = text
    return slide


def _add_text_slide(prs, title, lines):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = ln
        para.font.size = Pt(16)
    return slide


def _add_table_slide(prs, title, df, max_rows=12):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    d = df.head(max_rows).fillna("")
    rows, cols = d.shape
    tbl = slide.shapes.add_table(rows + 1, cols, Inches(0.4), Inches(1.4),
                                 Inches(12.5), Inches(0.4 * (rows + 1))).table
    for j, col in enumerate(d.columns):
        cell = tbl.cell(0, j)
        cell.text = str(col)
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].font.bold = True
    for i in range(rows):
        for j in range(cols):
            cell = tbl.cell(i + 1, j)
            cell.text = str(d.iat[i, j])
            cell.text_frame.paragraphs[0].font.size = Pt(10)
    return slide


def _add_image_slide(prs, title, image_path):
    from pptx.util import Inches
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(1.2), Inches(1.4), height=Inches(5.6))
    return slide


def _add_gallery_slide(prs, title, rep, own):
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    w = 2.0
    n = 0
    for shop_id, path in rep.items():
        p = ROOT / path
        if not p.exists():
            continue
        left = Inches(0.5 + (n % 6) * 2.1)
        top = Inches(1.5 + (n // 6) * 2.6)
        try:
            slide.shapes.add_picture(str(p), left, top, width=Inches(w))
        except Exception:
            continue
        lab = slide.shapes.add_textbox(left, Inches(1.5 + (n // 6) * 2.6 + 2.0),
                                       Inches(w), Inches(0.4))
        run = lab.text_frame.paragraphs[0]
        run.text = ("★ " if shop_id in own else "") + shop_id
        run.font.size = Pt(10)
        run.font.bold = shop_id in own
        n += 1
    if n == 0:  # 이미지 없으면 안내 문구
        box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(1))
        box.text_frame.text = "수집된 썸네일 없음 — collect.py 로 이미지 수집 후 재생성하세요."
    return slide


def add_marketing_slides(prs, group):
    """톤앤매너·메인 캠페인 / 인플루언서·IP 콜라보 / 이벤트·프로모션 슬라이드."""
    import marketing
    data = marketing.mine(group)

    # 1) 톤앤매너 + 메인 캠페인 (텍스트)
    lines = []
    for sid, d in data.items():
        star = "★ " if d["is_own"] else ""
        lines.append(f"{star}{sid} — {d['tone']}")
        lines.append(f"     ▸ {d['hero']}")
    _add_text_slide(prs, "마케팅 · 톤앤매너 & 메인 캠페인", lines)

    # 2) 인플루언서·IP 콜라보 (표)
    rows = []
    for sid, d in data.items():
        star = "★" if d["is_own"] else ""
        for ip in d["ip_collabs"]:
            rows.append({"샵": f"{star}{sid}", "파트너": ip, "유형": "IP", "상품수": "-"})
        for c in d["collabs"]:
            rows.append({"샵": f"{star}{sid}", "파트너": c["partner"],
                         "유형": "인플루언서/IP", "상품수": c["count"]})
    if rows:
        _add_table_slide(prs, "마케팅 · 인플루언서/IP 콜라보 리스트", pd.DataFrame(rows), max_rows=16)

    # 3) 이벤트·프로모션 현황 (표)
    prows = []
    for sid, d in data.items():
        star = "★" if d["is_own"] else ""
        s = d["snap"] or [None] * 8
        prows.append({"샵": f"{star}{sid}", "쿠폰": s[3], "최대할인": s[4],
                      "타임세일": s[5], "평균%↓": s[6], "메가와리": s[7],
                      "限定": d["kw"]["限定"], "セット": d["kw"]["セット"],
                      "先行": d["kw"]["先行"]})
    _add_table_slide(prs, "마케팅 · 이벤트/프로모션 현황", pd.DataFrame(prows), max_rows=16)


def add_group_slides(prs, group):
    conn = connect()
    own = own_shop_ids(group)
    df = build_scorecard(group)

    _section_title(prs, f"{group_label(group)}  ·  {group}")

    # 요약
    n_shops = len(df)
    leader = df.iloc[0]["shop_id"] if not df.empty and "review_volume" in df else "-"
    summary_lines = [
        f"· 분석 샵: {n_shops}개 (자사 {len(own)} / 경쟁사 {n_shops - len(own)})",
        f"· 리뷰볼륨 선두: {leader}",
        f"· 자사: {', '.join(own) or '-'}",
    ]
    _add_text_slide(prs, "요약", summary_lines)

    # 스코어카드
    if not df.empty:
        _add_table_slide(prs, "스코어카드", df)

    # 포지셔닝 맵
    for p in positioning_maps(df, group, own):
        _add_image_slide(prs, p.stem.replace(f"positioning_{group}_", "포지셔닝 · "), p)

    # 갤러리
    _add_gallery_slide(prs, "디자인 벤치마크 갤러리 (★ 자사)", gallery_assets(conn, group, own), own)

    # 마케팅 인텔리전스 (톤앤매너·콜라보·이벤트)
    add_marketing_slides(prs, group)

    # 갭 → 개선 과제
    gaps = gap_table(df, own)
    if not gaps.empty:
        _add_table_slide(prs, "갭 분석 (자사 − 그룹 1위, 갭 작은 순 = 우선순위)", gaps)
        tasks = [f"· [{r['지표']}] {r['그룹1위']} 대비 갭 {r['갭']} "
                 f"(달성률 {r['달성률%']}%) → 개선 검토"
                 for _, r in gaps[gaps["갭"] < 0].iterrows()]
        if tasks:
            _add_text_slide(prs, "개선 과제 (갭 음수 지표)", tasks)

    # 부록: 원본 스코어 데이터
    if not df.empty:
        _add_table_slide(prs, "부록 · 원본 스코어 데이터", df, max_rows=20)


# ── 산출물 ──────────────────────────────────────────────────────────────────
def generate_static(group):
    """PPTX 없이 맵 png + 갭 csv 만 생성."""
    df = build_scorecard(group)
    own = own_shop_ids(group)
    if df.empty:
        print(f"[{group}] 데이터 없음 — collect.py/score.py 먼저 실행.")
        return
    OUT.mkdir(exist_ok=True)
    maps = positioning_maps(df, group, own)
    for m in maps:
        print(f"  맵 저장: {m}")
    gaps = gap_table(df, own)
    if not gaps.empty:
        gpath = OUT / f"gap_{group}.csv"
        gaps.to_csv(gpath, index=False, encoding="utf-8-sig")
        print(f"  갭표 저장: {gpath}")


def main():
    want_pptx = "--pptx" in sys.argv
    merge = "--merge" in sys.argv
    group = _arg("--group", "skincare")

    if want_pptx:
        OUT.mkdir(exist_ok=True)
        if merge:
            prs = _new_pptx()
            _title_slide(prs, "큐텐 경쟁사 비교분석 (통합)",
                         f"그룹: {', '.join(GROUPS)} · 2026")
            for g in GROUPS:
                add_group_slides(prs, g)
            path = OUT / "report_merged.pptx"
        else:
            prs = _new_pptx()
            _title_slide(prs, f"큐텐 경쟁사 비교분석 — {group_label(group)}",
                         f"그룹: {group} · 자사 vs 경쟁사")
            add_group_slides(prs, group)
            path = OUT / f"report_{group}.pptx"
        prs.save(str(path))
        print(f"PPTX 저장: {path}")
    else:
        groups = list(GROUPS) if (group == "all" and merge) else [group]
        for g in groups:
            print(f"=== {g} ===")
            generate_static(g)


if __name__ == "__main__":
    main()
