"""재사용 슬라이드 컴포넌트 — python-pptx 저수준 래퍼.

slide 좌표/치수는 모두 Inches. 색은 theme.RGBColor.
"""
from __future__ import annotations
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from . import theme as T

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def blank_slide(prs):
    """레이아웃 없는 빈 슬라이드."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, fill=None, line=None, line_w=0.75, rounded=False, radius=0.06):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, Inches(l), Inches(t), Inches(w), Inches(h))
    if rounded:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def soft_shadow(shape, blur=0.06, dist=0.03, alpha=82):
    """카드용 부드러운 외부 그림자(XML 직접). alpha=투명도(%)."""
    from pptx.oxml.ns import qn
    spPr = shape._element.spPr
    # 기존 effectLst 제거
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    sh = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(Emu(Inches(blur))), "dist": str(Emu(Inches(dist))),
        "dir": "5400000", "rotWithShape": "0"})
    clr = sh.makeelement(qn("a:srgbClr"), {"val": "1A1430"})
    al = clr.makeelement(qn("a:alpha"), {"val": str(int((100 - alpha) * 1000))})
    clr.append(al); sh.append(clr); eff.append(sh); spPr.append(eff)


def text(slide, l, t, w, h, paras, size=12, color=None, bold=False,
         align="left", anchor="top", font=None, line_spacing=None, space_after=2):
    """텍스트 박스. paras = str | [str|dict].
    dict 키: t(필수), size, color, bold, align, font, color_hex.
    """
    if color is None:
        color = T.INK
    if font is None:
        font = T.FONT
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = ANCHOR[anchor]
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    items = paras if isinstance(paras, list) else [paras]
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        d = {"t": it} if isinstance(it, str) else dict(it)
        p.alignment = ALIGN[d.get("align", align)]
        if line_spacing:
            p.line_spacing = line_spacing
        p.space_after = Pt(d.get("space_after", space_after))
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = d["t"]
        f = run.font
        f.size = Pt(d.get("size", size))
        f.bold = d.get("bold", bold)
        f.name = d.get("font", font)
        c = d.get("color", color)
        f.color.rgb = c
    return box


def page_bg(slide, color=None):
    rect(slide, 0, 0, 13.333, 7.5, fill=color or T.PAGE_BG)


def header(slide, eyebrow, title, page_no, accent_title_part=None):
    """콘텐츠 슬라이드 상단: 짙은 헤더바 + 핑크 액센트선 + eyebrow/제목/페이지번호."""
    rect(slide, 0, 0, 13.333, 1.18, fill=T.INK)
    rect(slide, 0, 1.18, 13.333, 0.06, fill=T.PINK)
    text(slide, 0.60, 0.20, 9.0, 0.30, eyebrow, size=11, color=T.MUTED, bold=True)
    # 제목(일부를 핑크로 강조 가능)
    if accent_title_part and accent_title_part in title:
        head, _, tail = title.partition(accent_title_part)
        box = slide.shapes.add_textbox(Inches(0.60), Inches(0.44), Inches(11.0), Inches(0.62))
        tf = box.text_frame; tf.word_wrap = True
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        p = tf.paragraphs[0]
        for seg, col in ((head, T.WHITE), (accent_title_part, T.PINK), (tail, T.WHITE)):
            if not seg:
                continue
            r = p.add_run(); r.text = seg
            r.font.size = Pt(26); r.font.bold = True; r.font.name = T.FONT_BOLD; r.font.color.rgb = col
    else:
        text(slide, 0.60, 0.44, 11.0, 0.62, title, size=26, color=T.WHITE, bold=True, anchor="middle")
    # 페이지 번호
    text(slide, 11.93, 0.36, 0.90, 0.50, f"{page_no:02d}", size=22, color=T.PINK, bold=True, align="right")
    rect(slide, 12.87, 0.34, 0.02, 0.50, fill=T.MUTED)


def footer(slide, category_label, prepared="Aiden Lab"):
    text(slide, 0.60, 7.16, 9.0, 0.25,
         f"{prepared} · 큐텐 경쟁사 비교분석 ({category_label}) · 2026.06",
         size=8.5, color=T.MUTED)
    text(slide, 9.73, 7.16, 3.0, 0.25, "Confidential", size=8.5, color=T.MUTED, align="right")


def card(slide, l, t, w, h, fill=None, accent=None, accent_w=0.10, radius=0.05,
         shadow=True, line=None, line_w=1.0):
    """둥근 흰 카드 + (선택)좌측 액센트 바."""
    c = rect(slide, l, t, w, h, fill=fill or T.CARD_BG, rounded=True, radius=radius,
             line=line, line_w=line_w)
    if shadow:
        soft_shadow(c)
    if accent:
        rect(slide, l, t, accent_w, h, fill=accent, rounded=True, radius=0.25)
    return c


def chip(slide, l, t, w, h, label, fill=None, color=None, size=11, bold=True):
    rect(slide, l, t, w, h, fill=fill or T.ZEBRA, rounded=True, radius=0.3)
    text(slide, l, t, w, h, label, size=size, color=color or T.INK, bold=bold,
         align="center", anchor="middle")


def circle_num(slide, l, t, d, n, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = fill or T.PINK
    s.line.fill.background(); s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(n)
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = T.WHITE; r.font.name = T.FONT_BOLD
    return s


def image_fit(slide, path, l, t, w, h):
    """이미지를 박스에 center-crop으로 채움(가능하면). 실패 시 회색 플레이스홀더."""
    import os
    if not path or not os.path.exists(path):
        ph = rect(slide, l, t, w, h, fill=T.ZEBRA, rounded=True, radius=0.04)
        text(slide, l, t, w, h, "이미지 없음", size=10, color=T.MUTED, align="center", anchor="middle")
        return ph
    try:
        from PIL import Image
        iw, ih = Image.open(path).size
        box_ratio = w / h
        img_ratio = iw / ih
        if img_ratio > box_ratio:      # 이미지가 더 넓음 → 좌우 크롭
            crop_w = int(ih * box_ratio)
            x0 = (iw - crop_w) // 2
            crop = (x0, 0, x0 + crop_w, ih)
        else:                          # 더 높음 → 상하 크롭
            crop_h = int(iw / box_ratio)
            y0 = (ih - crop_h) // 2
            crop = (0, y0, iw, y0 + crop_h)
        import tempfile, hashlib
        Image.open(path).crop(crop).save(
            tmp := os.path.join(tempfile.gettempdir(),
                                "deckimg_" + hashlib.md5((path + str(crop)).encode()).hexdigest() + ".jpg"),
            quality=88)
        return slide.shapes.add_picture(tmp, Inches(l), Inches(t), Inches(w), Inches(h))
    except Exception:
        return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))


# 별칭
image_inside = image_fit
